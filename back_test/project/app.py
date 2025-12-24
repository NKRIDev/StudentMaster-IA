##pip install flask flask-cors requests PyPDF2 python-docx python-pptx ollama
#pip install python-dotenv

from flask import Blueprint

from project.models import User
from . import db

# Importing flask library
from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import io

# JWT token
from flask_jwt_extended import jwt_required, get_jwt_identity

# Setting up the Ollama library for AI
from ollama import chat
from ollama import ChatResponse

#Requests
import requests
import json

# Document treatment
import PyPDF2
from docx import Document
from pptx import Presentation

#Load .env
from dotenv import load_dotenv

load_dotenv()
OLLAMA_API = os.getenv("OLLAMA_API_URL")
FLASH_HOST = os.getenv("FLASK_HOST")
FLASK_PORT = int(os.getenv("FLASK_PORT"))

# Application and patch CORS errors
main = Blueprint('main', __name__)
CORS(main)

# Allow extensions
ALLOWED_EXTENSIONS = {'.pdf', '.txt', '.docx', '.md', '.pptx'}

# text file to be processed
CURRENT_FILE_NAME = None
CURRENT_FILE_TEXT = None

#Extracts text from a file based on its extension
def extract_text_from_file(file):
    _, extension = os.path.splitext(file.filename)
    extension = extension.lower()

    file_content = file.read()
    
    try:
        # Simple text files
        if extension in ['.txt', '.md']:
            return file_content.decode('utf-8')
        
        # PDF
        elif extension == '.pdf':
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text = []
            for page in pdf_reader.pages:
                text.append(page.extract_text())
            return '\n'.join(text)
        
        # Word => .docx
        elif extension == '.docx':
            doc = Document(io.BytesIO(file_content))
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        
        # PowerPoint => .pptx
        elif extension == '.pptx':
            prs = Presentation(io.BytesIO(file_content))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        
        else:
            return "Extension not supported"
            
    except Exception as e:
        return f"Error during extraction : {str(e)}"

# Generating the document summary with Ollama
def summarize(text, model="llama3.1:8b"):
    
    prompt = f"""
Tu es un assistant expert en synthèse documentaire.

Ta tâche : produire un **résumé clair, concis et parfaitement structuré en Markdown**, destiné à être affiché dans une interface web sous forme de **cards** ou **accordéons**.

Respecte STRICTEMENT les règles suivantes :

### ✅ Règles de rédaction
- Résumer uniquement les informations importantes.
- Éviter toute répétition ou reformulation inutile.
- Utiliser exclusivement du **Markdown propre et lisible**.
- Aucune phrase d’introduction ou conclusion de ta part.
- Aucune mention du prompt, du modèle ou du texte original.
- Pas de phrases du style “Voici le résumé” → commence directement par le contenu.
- Style direct, clair, simple et bien organisé.
- Uniquement du Markdown propre.
- Utilise des **accordéons (`<details>` / `<summary>`)** pour les sections principales.
- Utilise des **cards (`<div class="card">...</div>`)** pour sous-sections.
- Mets les points clés en **listes à puces**.
- Les informations importantes peuvent être mises dans des **blockquote**.
- Paragraphes courts et synthétiques.
- Markdown prêt à afficher dans une interface web interactive.

### 🧱 Structure attendue
# Titre principal

## Section
<details>
  <summary>Titre de la section</summary>
  <div class="card">
    Paragraphe court.
    <ul>
      <li>Point clé 1</li>
      <li>Point clé 2</li>
    </ul>
  </div>
  <blockquote>Une citation ou point important.</blockquote>
</details>

## Autre Section
<details>
  <summary>Titre de la section</summary>
  ...
</details>

### 📘 Texte à résumer (extrait)
{text}

### ➤ Génère maintenant UNIQUEMENT le résumé en Markdown interactif prêt à afficher :
"""
    
    try:
        response = requests.post(
            OLLAMA_API,
            json={
                "model": model,
                "prompt": prompt,
                "stream": False
            },
            timeout=120
        )
        
        if response.status_code == 200:
            return response.json()['response']
        else:
            return f"API error: {response.status_code}"
    except Exception as e:
        return f"Connection error: {str(e)}"

# Generate flashcards linked to the document using Ollama
def flashcards_generator(text, model="llama3.1:8b"):
    prompt = f"""
Tu es un assistant expert en pédagogie et en synthèse documentaire.

❌ NE METS JAMAIS de texte autre que le JSON.
❌ Pas de phrases d’introduction, pas de conclusion, pas de numérotation ou explications.
✅ Retourne uniquement un tableau JSON strict, prêt à parser.

Ta tâche : générer entre 5 et 10 **flashcards pédagogiques** à partir du texte fourni.  
Chaque flashcard doit être un objet JSON respectant l'interface TypeScript suivante :

export interface Flashcard {{
  id: string
  question: string
  answer: string
}}

### ✅ Consignes strictes
- Génère **uniquement des flashcards**.
- Chaque flashcard doit avoir :
  - `id` → identifiant unique (UUID ou n'importe quelle chaîne unique)
  - `question` → une question claire et concise
  - `answer` → réponse courte et directe, basée sur le texte
- Pas de contenu hors du format JSON
- Entre 5 et 10 flashcards maximum
- Questions et réponses directement liées aux informations importantes du texte
- Évite toute phrase introductive ou explicative
- Retourne un tableau JSON complet prêt à parser dans React/TypeScript
- Je veux que tu retournes uniquement le tableau, pas autre chose !

### Format attendu
[
  {{
    "id": "1",
    "question": "Quelle est la définition de X ?",
    "answer": "X est ..."
  }},
  {{
    "id": "2",
    "question": "Quels sont les avantages de Y ?",
    "answer": "Y permet ..."
  }}
]

### Texte source à partir duquel générer les flashcards
{text}

### ➤ Génère maintenant **UNIQUEMENT** les flashcards en JSON :
"""
    
    try:
        response = requests.post(
            OLLAMA_API,
            json={
                "model": model,
                "prompt": prompt,
                "stream": False
            },
            timeout=120
        )
        
        if response.status_code == 200:
            content = response.json()['response']

            try:
                flashcards_data = json.loads(content)
                return flashcards_data
            except json.JSONDecodeError:
                return json.dumps({"error", "Invalid JSON from flashcards"})
        else:
            return f"API error: {response.status_code}"
    except Exception as e:
        return f"Connection error: {str(e)}"

# The question was generated based on the document with Ollama
def quiz_generator(text, model="llama3.1:8b"):
    prompt = f"""
Tu es un assistant expert en pédagogie et génération de quiz.

❌ NE METS JAMAIS de texte autre que le JSON.
❌ AUCUNE phrase d’introduction, explication, note ou code TypeScript.
❌ AUCUN commentaire ou texte en dehors du tableau JSON.

✅ Retourne uniquement un tableau JSON strict, prêt à parser.

Ta tâche : générer entre 5 et 10 questions de **quiz** basées sur le texte fourni.  
Chaque question doit respecter strictement l’interface suivante :

export interface QuizQuestion {{
  id: string,
  question: string,
  options: string[],
  correctAnswer: number,
  explanation: string
}}

### ✅ Consignes strictes
- Retourne **EXCLUSIVEMENT** un tableau JSON.
- Entre **5 et 10** questions.
- Chaque élément doit contenir :
  - `id` → chaîne unique (UUID ou texte unique)
  - `question` → une question claire et précise
  - `options` → liste de **4 propositions maximum**
  - `correctAnswer` → index (0–3) de la bonne réponse
  - `explanation` → explication courte et correcte basée sur le texte
- Aucune phrase en dehors du JSON.
- Pas de code TypeScript.
- Pas d’autres champs que ceux définis.
- Le JSON doit être valide, propre et directement utilisable.

### Format (uniquement structure, ne pas réutiliser) :
[
  {{
    "id": "1",
    "question": "Exemple ?",
    "options": ["A", "B", "C", "D"],
    "correctAnswer": 1,
    "explanation": "Explication ..."
  }}
]

### Texte source
{text}

### ➤ Génère maintenant UNIQUEMENT le tableau JSON :
"""
    
    try:
        response = requests.post(
            OLLAMA_API,
            json={
                "model": model,
                "prompt": prompt,
                "stream": False
            },
            timeout=120
        )
        
        if response.status_code == 200:
            content = response.json()['response']
            
            try:
                quiz_data = json.loads(content)
                return quiz_data
            except json.JSONDecodeError:
                return json.dumps({"error", "Invalid JSON from Quizz"})
        else:
            return f"API error: {response.status_code}"
    except Exception as e:
        return f"Connection error : {str(e)}"

# Data from the application's various features. 
# This data is created when the user uploads their file.
SUMMARIZE = None
FLASHCARDS = None
QUIZ = None

# Retrieve the document to be processed
@main.route("/upload", methods=["POST"])
def upload_file():
    global CURRENT_FILE_NAME
    global CURRENT_FILE_TEXT

    if "file" not in request.files:
        return jsonify({"error" : "No file received"}), 400
    
    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error" : "Invalid file name"}), 400
    
    # Checked extension file
    _, extension = os.path.splitext(file.filename)
    if extension.lower() not in ALLOWED_EXTENSIONS:
        return jsonify({
            "error": f"L'extension '{extension}' n'est pas autorisé."
        }), 400
    
    # Get content file in text format
    CURRENT_FILE_NAME = file.filename
    CURRENT_FILE_TEXT = extract_text_from_file(file)

    #Checked content file
    if CURRENT_FILE_TEXT.startswith("Error") or CURRENT_FILE_TEXT == "Extension not supported":
        return jsonify({"error": CURRENT_FILE_TEXT}), 500
    
    # Create features datas
    summarize_data = summarize(CURRENT_FILE_TEXT)
    flashcards_data = flashcards_generator(CURRENT_FILE_TEXT)
    quiz_data = quiz_generator(CURRENT_FILE_TEXT)

    return jsonify({"filename": file.filename, 
                    "summarize" : summarize_data, 
                    "flashcards" : flashcards_data,
                    "quiz": quiz_data})

# Create a summary of the sent file in markdown format 
@main.route("/summary")
def get_summary():
    global CURRENT_FILE_NAME
    global SUMMARIZE

    if SUMMARIZE is None:
        return jsonify({"error" : "No file was sent"}), 400

    return jsonify({"message": "ok", "filename" : CURRENT_FILE_NAME, "content" : SUMMARIZE})

# Send flashcards
@main.route("/flashcards")
def get_flashcards():
    global CURRENT_FILE_NAME
    global FLASHCARDS

    if FLASHCARDS is None:
        return jsonify({"error" : "No file was sent"}), 400
    
    # Transform text format into json
    flashcards_data = None
    try:
        flashcards_data = json.loads(FLASHCARDS)
    except json.JSONDecodeError:
        return json.dumps({"error", "Invalid JSON from flashcards"})
    
    return jsonify({"message": "ok", "filename" : CURRENT_FILE_NAME, "content" : flashcards_data})

# Send quiz
@main.route("/quiz")
def get_quiz():
    global CURRENT_FILE_NAME
    global QUIZ

    if QUIZ is None:
        return jsonify({"error" : "No file was sent"}), 400
    
    # Transform text format into json
    quiz_data = None
    try:
        quiz_data = json.loads(QUIZ)
    except json.JSONDecodeError:
        return json.dumps({"error", "Invalid JSON from quiz"})

    return jsonify({"message": "ok", "filename" : CURRENT_FILE_NAME, "content" : quiz_data})

# Auth test
@main.route("/")
def index():
    return "Index"

# Profile page
@main.route("/api/profile", methods=["GET"])
@jwt_required()
def profile():
    user_id = get_jwt_identity()
    user = User.query.get(user_id)

    return jsonify({
        "user": {
            "id": user.id,
            "email": user.email,
            "pseudo": user.pseudo
        }
    })

# Start flask server
if __name__ == "__main__":
    main.run(
        host=FLASK_HOST,
        port=FLASK_PORT,
        debug=True
    )